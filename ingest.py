#!/usr/bin/env python3
import os
import sys
import re
import uuid
import json
import io
import google.generativeai as genai
from openai import OpenAI
import psycopg2
from psycopg2.extras import execute_values
from PyPDF2 import PdfReader
from pptx import Presentation
from dotenv import load_dotenv
from canvasapi import Canvas
from datetime import datetime
import tiktoken
import argparse

# --- Configuration ---

# Load environment variables from .env file
load_dotenv()

API_URL = "https://ufl.instructure.com/"

# Canvas API key
API_KEY = os.getenv("CANVAS_API_KEY")

# Initialize a new Canvas object
canvas = Canvas(API_URL, API_KEY)

# Set your OpenAI API key and database URL via environment variables.
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
DATABASE_URL = os.getenv("DATABASE_URL")

if not OPENAI_API_KEY or not DATABASE_URL or not GEMINI_API_KEY:
    print("Please set the OPENAI_API_KEY, GEMINI_API_KEY and DATABASE_URL environment variables.")
    sys.exit(1)

openai_client = OpenAI(api_key=OPENAI_API_KEY)
genai.configure(api_key=GEMINI_API_KEY)

# Define the chunking prompt (as described in the article)
CHUNKING_PROMPT = (
    "OCR the following page into Markdown. Tables should be formatted as HTML. "
    "Do not sorround your output with triple backticks.\n\n"
    "Chunk the document into sections of roughly 250 - 1000 words. Our goal is \n"
    "to identify parts of the page with same semantic theme. These chunks will \n"
    "be embedded and used in a RAG pipeline. \n\n"
    "Surround the chunks with <chunk> </chunk> html tags."
)

# Set the token input limit to prevent cutoffs during the chunking process
GEMINI_INPUT_LIMIT = 7300


# --- Get Bytes from a file (for files from Canvas API) ---       

def file_from_bytes(file):
    bytes = file.get_contents(binary=True)
    return io.BytesIO(bytes)


# --- PDF Extraction ---

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from all pages of the PDF."""
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text.strip()

# --- PPTX Extraction ---

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from all slides of the PPTX, including speaker notes."""
    presentation = Presentation(pptx_path)
    full_text = []
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + "\n"
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_text += "Speaker Notes: " + slide.notes_slide.notes_text_frame.text + "\n"
        full_text.append(slide_text)
    return "\n".join(full_text).strip()


def split_text_into_safe_chunks(text: str) -> list:
    """Splits text into smaller chunks that fit within Gemini's token limit."""
    encoding = tiktoken.encoding_for_model("text-embedding-ada-002")
    tokens = encoding.encode(text)

    # Split into chunks of GEMINI_INPUT_LIMIT tokens
    chunks = [tokens[i:i + GEMINI_INPUT_LIMIT] for i in range(0, len(tokens), GEMINI_INPUT_LIMIT)]

    # Convert tokens back to text
    return [encoding.decode(chunk) for chunk in chunks]


# --- Gemini Flash 2.0 Chunking ---
def chunk_text_with_gemini(text: str) -> list:
    """
    Send the extracted text and chunking prompt to Gemini Flash 2.0.
    Expected output: a markdown string containing one or more <chunk>...</chunk> sections.
    """
    try:

        """Send smaller text chunks to Gemini to avoid truncation issues."""
        text_chunks = split_text_into_safe_chunks(text)  # Pre-split input
        all_chunks = []

        generation_config = {
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 40,
            "max_output_tokens": 8192,
            "response_mime_type": "text/plain",
        }

        model = genai.GenerativeModel(
            model_name="gemini-2.0-flash",
            generation_config=generation_config,
        )

        chat_session = model.start_chat(history=[])
        for chunk in text_chunks:
            response = chat_session.send_message(f"{CHUNKING_PROMPT}\n\n{text}")
            all_chunks.extend(response)

            # chunked_text = response.text
            # Save raw chunked text to file for debugging
            # with open("chunks_debug.txt", "w") as f:
            #     f.write(chunked_text)

        
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
    
    # Extract chunks enclosed in <chunk>...</chunk> tags using regex.
    chunks = []
    for chunk in all_chunks:
        chunks.extend(re.findall(r"<chunk>(.*?)</chunk>", chunk.text, re.DOTALL))
    if len(chunks) == 0:
        print("No chunks found in the Gemini response. Check the API output.")

    # Save extracted chunks to file for debugging
    # with open("extracted_chunks.txt", "w") as f:
    #     for i, chunk in enumerate(chunks):
    #         f.write(f"\n--- Chunk {i+1} ---\n")
    #         f.write(chunk.strip())
    #         f.write("\n")

    return [chunk.strip() for chunk in chunks if chunk.strip()]

# --- Embedding Generation ---

def generate_embedding(text: str) -> list:
    """Generate an embedding vector for a given text using OpenAI."""
    try:
        response = openai_client.embeddings.create(
            input=text.replace("\n", " "),
            model="text-embedding-ada-002"
        )
        # The API returns a list in response['data'][0]['embedding']
        return response.data[0].embedding
    except Exception as e:
        print(f"Error generating embedding: {e}")

# --- Database Insertion ---

def store_resource_and_embeddings(full_text: str, chunks: list, embeddings: list, filename: str = None, url: str = None, file_path: str = None, file_id: int = None):
    """
    Insert a new resource (course slide content) into the resources table,
    and then insert each chunk along with its embedding into the embeddings table.
    """
    resource_id = file_id
    # Prepare the embeddings rows as tuples: (id, resource_id, content, embedding)
    embedding_rows = []
    for chunk, emb in zip(chunks, embeddings):
        emb_id = str(uuid.uuid4())
        emb_array = "[" + ",".join(map(str, emb)) + "]"
        embedding_rows.append((emb_id, resource_id, chunk, emb_array))

    try:
        conn = psycopg2.connect(DATABASE_URL)
        conn.autocommit = True
        cur = conn.cursor()

        # Insert resource record with filename and url
        cur.execute(
            """
            INSERT INTO resources (id, content, filename, url, filepath)
            VALUES (%s, %s, %s, %s, %s)
            ON CONFLICT (id) DO UPDATE 
            SET content = EXCLUDED.content,
                filename = EXCLUDED.filename,
                url = EXCLUDED.url,
                filepath = EXCLUDED.filepath
            RETURNING (xmax <> 0) AS updated;
            """,
            (resource_id, full_text, filename, url, file_path)
        )

        updated = cur.fetchone()[0]  # True if an update occurred

        # If an update occurred, delete associated embeddings
        if updated:
            delete_query = "DELETE FROM embeddings WHERE resource_id = %s;"
            cur.execute(delete_query, (str(resource_id),))

        # Bulk insert into the embeddings table
        insert_query = """
            INSERT INTO embeddings (id, resource_id, content, embedding)
            VALUES %s
        """
        execute_values(cur, insert_query, embedding_rows)
        print("Resource and embeddings successfully stored. \n \n")
    except Exception as e:
        print(f"Database error: {e}")
    finally:
        if conn:
            cur.close()
            conn.close()

# --- Retrieve all ids for comparison with reingested files ---

def get_all_resource_ids():
    try:
        conn = psycopg2.connect(DATABASE_URL)
        conn.autocommit = True
        cur = conn.cursor()

        # Fetch all resource IDs
        cur.execute("SELECT id FROM resources")
        resource_ids = {int(row[0]) for row in cur.fetchall()}  # Extract IDs into a set

        return resource_ids
    
    except Exception as e:
        print(f"Database error: {e}")

    finally:
        if conn:
            # Close the connection
            cur.close()
            conn.close()

# --- Delete a Resource and its Corresponding Embeddings ---          

def delete_file_from_index(file_id: int):
    try:
        conn = psycopg2.connect(DATABASE_URL)
        conn.autocommit = True
        cur = conn.cursor()

        cur.execute("DELETE FROM resources WHERE id = %s", (str(file_id),))

    except Exception as e:
        print(f"Database error: {e}")

    finally:
        cur.close()
        conn.close()


# --- Main Ingestion Pipeline ---

def main():
    parser = argparse.ArgumentParser(description='Ingest files into the knowledge base.')
    parser.add_argument('-r', '--reingest', action='store_true', help="Flag to trigger re-ingestion")
    args = parser.parse_args()

    # Get course by course number
    course = canvas.get_course(525691)

    file_ids = set()

    files = course.get_files()

    resource_ids = get_all_resource_ids()

    # iterate through files in Canvas course
    for file in files:

        # keep track of all current files in the canvas course
        file_ids.add(file.id)

        # if reingesting, check if the file was updated in the past 24 hours, if not continue
        if args.reingest:
            current_time = datetime.now()
            updated_at = datetime.fromisoformat(file.updated_at.rstrip("Z"))

            # second condition in if statement in case some file was deleted and never re-added in previous reingestion
            if (current_time - updated_at).days < 1 or file.id not in resource_ids:
                print("File was updated in the last 24 hours or is not yet ingested.")
                print("Date last modified:", updated_at)
            else:
                print("Skipping", file.display_name, "not updated. \n")
                continue


        filename = file.display_name
        file_path = course.get_folder(file.folder_id).full_name + "/" + filename
        file_path = file_path[7:] # remove "course " from file path
        url = file.url.split('/download?download_frd')[0]
        _, ext = os.path.splitext(filename)
        full_text = ''

        if ext == ".pdf":
            file_like = file_from_bytes(file)
            full_text = extract_text_from_pdf(file_like)
        elif ext == ".pptx":
            file_like = file_from_bytes(file)
            full_text = extract_text_from_pptx(file_like)
        else:
            print("Not ingesting file type:", ext)
            print()
            continue

        print("File being ingested:", filename)
        print()

        # eliminate null bytes:
        if "\x00" in full_text:
            print("Warning: Null byte detected in full_text!")
            full_text = full_text.replace("\x00", "")

        # chunk text
        print("Chunking text with Gemini Flash 2.0 ...")
        chunks = chunk_text_with_gemini(full_text)
        print(f"Found {len(chunks)} chunks.")

        # generate embeddings
        print("Generating embeddings for each chunk ...")
        embeddings = []
        for chunk in chunks:
            emb = generate_embedding(chunk)
            embeddings.append(emb)
        
        # store in db
        print("Storing resource and embeddings into the database ...")
        store_resource_and_embeddings(full_text, chunks, embeddings, filename, url, file_path, file.id)

    for resource_id in resource_ids:
        if resource_id not in file_ids:
            print("Deleting old file. File ID:", resource_id)
            delete_file_from_index(resource_id)

   

if __name__ == "__main__":
    main()
